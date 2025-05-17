# Tkinter Desktop App

import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches
import openai
import os
import requests
from io import BytesIO
import threading
import pandas as pd

openai.api_key = "YOUR_OPENAI_API_KEY"

def generate_slide_content(topic):
    prompt = f"""
Create a PowerPoint presentation on the topic: "{topic}".
Structure output exactly like this:

Slide 1 Title: <title>
Slide 1 Content: <bullet1>\n<bullet2>\n<bullet3>

Slide 2 Title: ...
..."""
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    return response['choices'][0]['message']['content']

def generate_image(prompt):
    response = openai.Image.create(prompt=prompt, n=1, size="512x512")
    image_url = response['data'][0]['url']
    image_data = requests.get(image_url).content
    return BytesIO(image_data)

def parse_and_create_ppt(text, topic):
    prs = Presentation()
    slides_data = text.strip().split("\n\n")

    for i, slide in enumerate(slides_data):
        lines = slide.strip().split("\n")
        title = lines[0].split(": ", 1)[1]
        bullets = [line.split(": ", 1)[1] if ": " in line else line for line in lines[1:]]
        content_text = '\n'.join(bullets)

        slide_layout = prs.slide_layouts[1]
        slide_obj = prs.slides.add_slide(slide_layout)
        slide_obj.shapes.title.text = title
        slide_obj.placeholders[1].text = content_text

        try:
            img = generate_image(f"{title}, {bullets[0]}")
            left = Inches(5.5)
            top = Inches(1.5)
            height = Inches(3.5)
            slide_obj.shapes.add_picture(img, left, top, height=height)
        except Exception as e:
            print(f"Image skipped: {e}")

    os.makedirs("gui_presentations", exist_ok=True)
    output_path = os.path.join("gui_presentations", f"{topic[:50].replace(' ', '_')}.pptx")
    prs.save(output_path)
    return output_path

def generate_from_topic(topic, log_box):
    log_box.insert(tk.END, f"\nGenerating for topic: {topic}\n")
    content = generate_slide_content(topic)
    ppt_path = parse_and_create_ppt(content, topic)
    log_box.insert(tk.END, f"Saved: {ppt_path}\n")

def run_single():
    topic = entry.get()
    if not topic:
        messagebox.showerror("Error", "Please enter a topic")
        return
    threading.Thread(target=generate_from_topic, args=(topic, log_box)).start()

def run_batch():
    file_path = filedialog.askopenfilename(filetypes=[("CSV or TXT", "*.csv *.txt")])
    if not file_path:
        return
    if file_path.endswith(".txt"):
        with open(file_path, "r") as f:
            topics = [line.strip() for line in f if line.strip()]
    elif file_path.endswith(".csv"):
        df = pd.read_csv(file_path)
        if 'topic' not in df.columns:
            messagebox.showerror("Error", "CSV must have 'topic' column")
            return
        topics = df['topic'].dropna().tolist()
    else:
        messagebox.showerror("Error", "Invalid file format")
        return

    for topic in topics:
        threading.Thread(target=generate_from_topic, args=(topic, log_box)).start()

# === GUI ===
root = tk.Tk()
root.title("Slide Generator AI")
root.geometry("600x500")

frame = tk.Frame(root)
frame.pack(pady=10)

entry = tk.Entry(frame, width=50)
entry.pack(side=tk.LEFT, padx=5)
btn_generate = tk.Button(frame, text="Generate Slide (Single Topic)", command=run_single)
btn_generate.pack(side=tk.LEFT)

btn_batch = tk.Button(root, text="Batch Upload (TXT/CSV)", command=run_batch)
btn_batch.pack(pady=10)

log_box = tk.Text(root, height=20, width=70)
log_box.pack()

root.mainloop()
