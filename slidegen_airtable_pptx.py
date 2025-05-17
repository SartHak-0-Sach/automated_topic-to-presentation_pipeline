# Airtable Scheduler that uses python-pptx instead of Google Slides

import requests
import openai
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from datetime import datetime
import smtplib
from email.message import EmailMessage
import os
import logging
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE
from flask import Flask, render_template_string, send_from_directory
import threading
import webbrowser

# CONFIGURATION
AIRTABLE_API_KEY = 'YOUR_AIRTABLE_API_KEY'
AIRTABLE_BASE_ID = 'YOUR_BASE_ID'
AIRTABLE_TABLE_NAME = 'Topics'
OPENAI_API_KEY = 'YOUR_OPENAI_API_KEY'
EMAIL_SENDER = 'you@example.com'
EMAIL_RECEIVER = 'receiver@example.com'
EMAIL_PASSWORD = 'your_email_password'
LOG_FILE = 'slidegen.log'
SLIDE_FOLDER = 'slides'
os.makedirs(SLIDE_FOLDER, exist_ok=True)
openai.api_key = OPENAI_API_KEY

# Setup logging
logging.basicConfig(
    filename=LOG_FILE,
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

# Airtable URL
AIRTABLE_ENDPOINT = f"https://api.airtable.com/v0/{AIRTABLE_BASE_ID}/{AIRTABLE_TABLE_NAME}"
HEADERS = {'Authorization': f'Bearer {AIRTABLE_API_KEY}'}

app = Flask(__name__)

@app.route('/')
def dashboard():
    with open(LOG_FILE, 'r') as f:
        logs = f.read().split('\n')[-100:]
    files = os.listdir(SLIDE_FOLDER)
    return render_template_string("""
    <h2>Slidegen Dashboard</h2>
    <h3>Generated Slides</h3>
    <ul>{% for file in files %}<li><a href='/slides/{{file}}'>{{file}}</a></li>{% endfor %}</ul>
    <h3>Latest Logs</h3>
    <pre>{{ logs }}</pre>
    """, logs='\n'.join(logs), files=files)

@app.route('/slides/<path:filename>')
def download_file(filename):
    return send_from_directory(SLIDE_FOLDER, filename, as_attachment=True)

def run_dashboard():
    threading.Thread(target=lambda: app.run(debug=False, port=7860)).start()
    time.sleep(1)
    webbrowser.open("http://localhost:7860")

def fetch_topics():
    response = requests.get(AIRTABLE_ENDPOINT, headers=HEADERS)
    records = response.json().get('records', [])
    topics = []
    for record in records:
        topic = record['fields'].get('Topic')
        status = record['fields'].get('Status', '')
        if topic and status != 'Done':
            topics.append((record['id'], topic))
    return topics

def mark_topic_done(record_id):
    url = f"{AIRTABLE_ENDPOINT}/{record_id}"
    data = {"fields": {"Status": "Done"}}
    requests.patch(url, headers={**HEADERS, "Content-Type": "application/json"}, json=data)

def generate_slide_content(topic):
    prompt = f"""
Create a PowerPoint presentation on the topic: "{topic}".
Structure the output:
Slide 1 Title: <title>
Slide 1 Content: <bullet1>\n<bullet2>\n<bullet3>
..."""
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    return response['choices'][0]['message']['content']

def generate_image(prompt):
    try:
        response = openai.Image.create(
            prompt=prompt,
            n=1,
            size="512x512"
        )
        image_url = response['data'][0]['url']
        image_data = requests.get(image_url).content
        return BytesIO(image_data)
    except Exception as e:
        logging.error(f"Image generation failed for prompt '{prompt}': {e}")
        return None

def parse_and_create_ppt(text, topic):
    prs = Presentation()
    slides_data = text.strip().split("\n\n")

    for i, slide in enumerate(slides_data):
        lines = slide.strip().split("\n")
        title = lines[0].split(": ", 1)[1] if ": " in lines[0] else f"Slide {i+1}"
        bullets = [line.split(": ", 1)[1] if ": " in line else line for line in lines[1:]]

        slide_layout = prs.slide_layouts[5]  # blank layout
        slide_obj = prs.slides.add_slide(slide_layout)

        title_shape = slide_obj.shapes.title or slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(0.2), Inches(8), Inches(1))
        title_shape.text = title

        content_box = slide_obj.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(4))
        tf = content_box.text_frame
        tf.word_wrap = True
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)

        img_stream = generate_image(f"{title} - {topic}")
        if img_stream:
            slide_obj.shapes.add_picture(img_stream, Inches(5.5), Inches(1.5), width=Inches(4))

    filename = f"{topic[:40].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    path = os.path.join(SLIDE_FOLDER, filename)
    prs.save(path)
    return path

def send_email_with_attachment(subject, filename):
    msg = EmailMessage()
    msg["From"] = EMAIL_SENDER
    msg["To"] = EMAIL_RECEIVER
    msg["Subject"] = subject
    msg.set_content(f"Your AI-generated presentation on '{subject}' is attached.")

    with open(filename, 'rb') as f:
        msg.add_attachment(f.read(), maintype='application', subtype='vnd.openxmlformats-officedocument.presentationml.presentation', filename=os.path.basename(filename))

    with smtplib.SMTP_SSL('smtp.gmail.com', 465) as smtp:
        smtp.login(EMAIL_SENDER, EMAIL_PASSWORD)
        smtp.send_message(msg)

def batch_generate():
    logging.info("Checking Airtable for new topics...")
    topics = fetch_topics()
    if not topics:
        logging.info("No new topics found.")
    else:
        for record_id, topic in topics:
            try:
                logging.info(f"Generating slide deck for: {topic}")
                content = generate_slide_content(topic)
                pptx_file = parse_and_create_ppt(content, topic)
                send_email_with_attachment(topic, pptx_file)
                mark_topic_done(record_id)
                logging.info(f"Successfully created and emailed deck for: {topic}")
                time.sleep(2)
            except Exception as e:
                logging.error(f"Error processing topic '{topic}': {e}")

if __name__ == "__main__":
    run_dashboard()
    batch_generate()
