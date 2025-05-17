# GUI Version 2: Streamlit Web App

import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import openai
import os
import requests
from io import BytesIO
from datetime import datetime

openai.api_key = "YOUR_OPENAI_API_KEY"

def generate_slide_content(topic):
    prompt = f"""
Create a PowerPoint presentation on the topic: "{topic}".
Structure output exactly like this:

Slide 1 Title: <title>
Slide 1 Content: <bullet1>\n<bullet2>\n<bullet3>

Slide 2 Title: ...
..."""
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    return response['choices'][0]['message']['content']

def generate_image(prompt):
    response = openai.Image.create(prompt=prompt, n=1, size="512x512")
    image_url = response['data'][0]['url']
    image_data = requests.get(image_url).content
    return BytesIO(image_data)

def parse_and_create_ppt(text, topic):
    prs = Presentation()
    slides_data = text.strip().split("\n\n")

    for i, slide in enumerate(slides_data):
        lines = slide.strip().split("\n")
        title = lines[0].split(": ", 1)[1]
        bullets = [line.split(": ", 1)[1] if ": " in line else line for line in lines[1:]]
        content_text = '\n'.join(bullets)

        slide_layout = prs.slide_layouts[1]
        slide_obj = prs.slides.add_slide(slide_layout)
        slide_obj.shapes.title.text = title
        slide_obj.placeholders[1].text = content_text

        try:
            img = generate_image(f"{title}, {bullets[0]}")
            left = Inches(5.5)
            top = Inches(1.5)
            height = Inches(3.5)
            slide_obj.shapes.add_picture(img, left, top, height=height)
        except Exception as e:
            print(f"Image skipped: {e}")

    output_path = f"{topic[:40].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(output_path)
    return output_path

# === Streamlit UI ===
st.set_page_config(page_title="Slide Generator AI", layout="centered")
st.title("📊 AI Slide Deck Generator")

topic = st.text_input("Enter a Topic")

if st.button("Generate Presentation"):
    if not topic:
        st.warning("Please enter a topic.")
    else:
        with st.spinner("Generating content and images..."):
            content = generate_slide_content(topic)
            ppt_path = parse_and_create_ppt(content, topic)

        with open(ppt_path, "rb") as f:
            st.success("Presentation generated!")
            st.download_button("📥 Download PPTX", f, file_name=os.path.basename(ppt_path))
